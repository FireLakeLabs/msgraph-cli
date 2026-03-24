#!/usr/bin/env python3
"""Generate test fixture files for Phase 4 validation.

Creates:
  tests/fixtures/test-workbook.xlsx — basic workbook with Sheet1 and sample data
  tests/fixtures/append-test.xlsx   — workbook with Sheet1 + Table1 (3 columns)
  tests/fixtures/test-doc.docx      — simple Word document with text
  tests/fixtures/test-slide.pptx    — single-slide presentation with text

Prerequisites:
  pip install openpyxl python-docx python-pptx
"""
from pathlib import Path


def create_basic_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Score", "Grade"])
    ws.append(["Alice", 95, "A"])
    ws.append(["Bob", 88, "B+"])
    ws.append(["Charlie", 72, "C"])
    wb.save(path)
    print(f"Created {path}")


def create_append_xlsx(path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.worksheet.table import Table, TableStyleInfo

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Score", "Grade"])
    ws.append(["Alice", 95, "A"])
    tab = Table(displayName="Table1", ref="A1:C2")
    tab.tableStyleInfo = TableStyleInfo(
        name="TableStyleMedium9",
        showFirstColumn=False,
        showLastColumn=False,
        showRowStripes=True,
    )
    ws.add_table(tab)
    wb.save(path)
    print(f"Created {path}")


def create_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Phase 4 Test Document", level=1)
    doc.add_paragraph("This is a test document for msgraph-cli validation.")
    doc.add_paragraph("It contains simple text content for the docs cat command.")
    doc.save(path)
    print(f"Created {path}")


def create_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Phase 4 Test"
    slide.placeholders[1].text = "This is a test slide for msgraph-cli validation."
    prs.save(path)
    print(f"Created {path}")


if __name__ == "__main__":
    out = Path(__file__).resolve().parent.parent / "tests" / "fixtures"
    out.mkdir(parents=True, exist_ok=True)

    create_basic_xlsx(out / "test-workbook.xlsx")
    create_append_xlsx(out / "append-test.xlsx")
    create_docx(out / "test-doc.docx")
    create_pptx(out / "test-slide.pptx")
